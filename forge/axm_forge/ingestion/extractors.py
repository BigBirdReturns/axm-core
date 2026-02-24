"""
AXM Forge — Universal Document Extractors (Topology-Preserving)

INVARIANT: Documents are never flattened into a single string.
Every extractor yields DocumentBlocks bound to their structural position.

Structured formats (CSV, JSON) yield tier0_candidates that bypass the LLM.

See: IDENTITY.md for ID generation rules.
See: EXTENSIONS_REGISTRY.md for ext/locators@1 schema.
"""
from __future__ import annotations

import csv
import io
import json
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional


# ---------------------------------------------------------------------------
# Data model — plain dicts for Locators, no phantom imports
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class DocumentBlock:
    """A segment of text permanently bound to its structural position.

    locator is a plain dict, e.g.:
        {"kind": "pdf", "page_index": 0, "file_path": "/path/to/doc.pdf"}
        {"kind": "docx", "paragraph_index": 4, "file_path": "/path/to/doc.docx"}
        {"kind": "txt", "file_path": "/path/to/doc.txt"}

    The locator dict flows through the pipeline:
        extractor -> segmenter -> binder -> candidates.jsonl -> compiler -> ext/locators.parquet
    """
    text: str
    locator: Dict[str, Any]


@dataclass
class ExtractedDocument:
    """Result of ingestion.

    For unstructured documents: blocks is populated. Each block carries a Locator.
    For structured data: tier0_candidates is populated. Blocks may be empty.
    Both can coexist (e.g., a spreadsheet with headers and structured rows).
    """
    source_path: str
    format: str
    blocks: List[DocumentBlock] = field(default_factory=list)
    tier0_candidates: Optional[List[Dict[str, Any]]] = None
    page_count: Optional[int] = None
    metadata: Dict[str, Any] = field(default_factory=dict)

    @property
    def full_text(self) -> str:
        """Concatenated text from all blocks. Use ONLY for writing source.txt.
        The Locator data is preserved separately in the pipeline."""
        return "\n\n".join(b.text for b in self.blocks)


# ---------------------------------------------------------------------------
# Extractors — each preserves structural topology
# ---------------------------------------------------------------------------

def extract_txt(path: Path) -> ExtractedDocument:
    text = path.read_text(encoding="utf-8", errors="replace")
    block = DocumentBlock(
        text=text,
        locator={"kind": "txt", "file_path": str(path)},
    )
    return ExtractedDocument(blocks=[block], source_path=str(path), format="txt")


def extract_markdown(path: Path) -> ExtractedDocument:
    text = path.read_text(encoding="utf-8", errors="replace")
    block = DocumentBlock(
        text=text,
        locator={"kind": "markdown", "file_path": str(path)},
    )
    return ExtractedDocument(blocks=[block], source_path=str(path), format="markdown")


def extract_html(path: Path) -> ExtractedDocument:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        blocks = []
        for i, el in enumerate(soup.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6", "li", "td", "pre"])):
            text = el.get_text(strip=True)
            if text:
                blocks.append(DocumentBlock(
                    text=text,
                    locator={"kind": "html", "element_index": i, "tag": el.name, "file_path": str(path)},
                ))

        if not blocks:
            full = soup.get_text(separator="\n", strip=True)
            blocks = [DocumentBlock(text=full, locator={"kind": "html", "file_path": str(path)})]

    except ImportError:
        text = re.sub(r"<[^>]+>", " ", raw)
        text = re.sub(r"\s+", " ", text).strip()
        blocks = [DocumentBlock(text=text, locator={"kind": "html", "file_path": str(path)})]

    return ExtractedDocument(blocks=blocks, source_path=str(path), format="html")


def extract_pdf(path: Path) -> ExtractedDocument:
    """One block per page. Page index preserved in Locator."""
    blocks: List[DocumentBlock] = []

    try:
        import fitz
        doc = fitz.open(str(path))
        for page_num, page in enumerate(doc):
            text = page.get_text().strip()
            if text:
                blocks.append(DocumentBlock(
                    text=text,
                    locator={"kind": "pdf", "page_index": page_num, "file_path": str(path)},
                ))
        return ExtractedDocument(blocks=blocks, source_path=str(path), format="pdf", page_count=len(doc))
    except ImportError:
        pass

    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                text = (page.extract_text() or "").strip()
                if text:
                    blocks.append(DocumentBlock(
                        text=text,
                        locator={"kind": "pdf", "page_index": page_num, "file_path": str(path)},
                    ))
            return ExtractedDocument(blocks=blocks, source_path=str(path), format="pdf", page_count=len(pdf.pages))
    except ImportError:
        pass

    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        for page_num, page in enumerate(reader.pages):
            text = (page.extract_text() or "").strip()
            if text:
                blocks.append(DocumentBlock(
                    text=text,
                    locator={"kind": "pdf", "page_index": page_num, "file_path": str(path)},
                ))
        return ExtractedDocument(blocks=blocks, source_path=str(path), format="pdf", page_count=len(reader.pages))
    except ImportError:
        pass

    raise ImportError(
        "No PDF library available. Install one of: pymupdf, pdfplumber, PyPDF2\n"
        "  pip install pymupdf --break-system-packages"
    )


def extract_docx(path: Path) -> ExtractedDocument:
    """One block per paragraph. Paragraph index preserved in Locator."""
    try:
        from docx import Document
        doc = Document(str(path))
        blocks = []
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if text:
                blocks.append(DocumentBlock(
                    text=text,
                    locator={"kind": "docx", "paragraph_index": i, "file_path": str(path)},
                ))
        return ExtractedDocument(blocks=blocks, source_path=str(path), format="docx")
    except ImportError:
        raise ImportError("python-docx not available. Install: pip install python-docx --break-system-packages")


def extract_xlsx(path: Path) -> ExtractedDocument:
    """One block per sheet. Also emits tier0_candidates for structured rows."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True)
        blocks = []
        candidates = []

        for sheet in wb.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue

            lines = []
            for row in rows:
                line = "\t".join(str(c) if c is not None else "" for c in row)
                if line.strip():
                    lines.append(line)

            if lines:
                blocks.append(DocumentBlock(
                    text="\n".join(lines),
                    locator={"kind": "xlsx", "sheet_name": sheet.title, "file_path": str(path)},
                ))

            headers = [str(c).strip() if c else f"col_{j}" for j, c in enumerate(rows[0])]
            for row_idx, row in enumerate(rows[1:], start=1):
                if not any(row):
                    continue
                subject = str(row[0]).strip() if row[0] else f"record_{row_idx}"
                for col_idx, val in enumerate(row):
                    if col_idx == 0 or val is None:
                        continue
                    val_str = str(val).strip()
                    if val_str and col_idx < len(headers):
                        candidates.append({
                            "subject": subject,
                            "predicate": headers[col_idx],
                            "object": val_str,
                            "object_type": "literal:string",
                            "tier": 0,
                            "evidence": f"{headers[col_idx]}: {val_str}",
                            "locator": {"kind": "xlsx", "sheet_name": sheet.title,
                                        "row_index": row_idx, "file_path": str(path)},
                        })

        return ExtractedDocument(
            blocks=blocks, source_path=str(path), format="xlsx",
            tier0_candidates=candidates if candidates else None,
            metadata={"sheet_count": len(wb.worksheets)},
        )
    except ImportError:
        raise ImportError("openpyxl not available. Install: pip install openpyxl --break-system-packages")


def extract_pptx(path: Path) -> ExtractedDocument:
    """One block per slide. Slide index preserved in Locator."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        blocks = []
        for i, slide in enumerate(prs.slides):
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                blocks.append(DocumentBlock(
                    text="\n".join(texts),
                    locator={"kind": "pptx", "slide_index": i, "file_path": str(path)},
                ))
        return ExtractedDocument(blocks=blocks, source_path=str(path), format="pptx", page_count=len(prs.slides))
    except ImportError:
        raise ImportError("python-pptx not available. Install: pip install python-pptx --break-system-packages")


def extract_csv(path: Path) -> ExtractedDocument:
    """Schema IS extraction. Bypasses LLM entirely.

    Emits tier0_candidates with tier=0.
    subject = first column, predicate = column header, object = cell value.
    """
    raw = path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(raw))
    rows = list(reader)

    if not rows:
        return ExtractedDocument(source_path=str(path), format="csv")

    headers = [h.strip() for h in rows[0]]
    candidates = []

    for row_idx, row in enumerate(rows[1:], start=1):
        if not any(row):
            continue
        subject = row[0].strip() if row else f"record_{row_idx}"
        for col_idx, val in enumerate(row):
            val = val.strip() if val else ""
            if col_idx > 0 and col_idx < len(headers) and val:
                candidates.append({
                    "subject": subject,
                    "predicate": headers[col_idx],
                    "object": val,
                    "object_type": "literal:string",
                    "tier": 0,
                    "evidence": f"{headers[col_idx]}: {val}",
                    "locator": {"kind": "csv", "row_index": row_idx,
                                "column": headers[col_idx], "file_path": str(path)},
                })

    block_text = "\n".join("\t".join(row) for row in rows)
    blocks = [DocumentBlock(text=block_text, locator={"kind": "csv", "file_path": str(path)})]

    return ExtractedDocument(
        blocks=blocks, source_path=str(path), format="csv",
        tier0_candidates=candidates if candidates else None,
        metadata={"headers": headers, "row_count": len(rows) - 1},
    )


def extract_json(path: Path) -> ExtractedDocument:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        data = json.loads(raw)
        text = json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        text = raw
    block = DocumentBlock(text=text, locator={"kind": "json", "file_path": str(path)})
    return ExtractedDocument(blocks=[block], source_path=str(path), format="json")


def extract_xml(path: Path) -> ExtractedDocument:
    raw = path.read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"<[^>]+>", " ", raw)
    text = re.sub(r"\s+", " ", text).strip()
    block = DocumentBlock(text=text, locator={"kind": "xml", "file_path": str(path)})
    return ExtractedDocument(blocks=[block], source_path=str(path), format="xml")


def extract_xbrl(path: Path) -> ExtractedDocument:
    """XBRL financial reports — schema IS extraction, tier-0 candidates bypass LLM.

    Parses contextRef facts from XBRL XML into structured candidates.
    Subject = filing entity, predicate = XBRL concept, object = value.
    """
    import xml.etree.ElementTree as ET

    raw = path.read_text(encoding="utf-8", errors="replace")
    candidates = []

    try:
        root = ET.fromstring(raw)
    except ET.ParseError:
        # Fallback: treat as plain XML text
        return extract_xml(path)

    # Find entity identifier for subject
    entity_label = "entity"
    for elem in root.iter():
        if elem.tag.endswith("identifier") or elem.tag.endswith("EntityCommonStockSharesOutstanding"):
            if elem.text and elem.text.strip():
                entity_label = elem.text.strip()[:80]
                break

    for elem in root.iter():
        if elem.get("contextRef") is None:
            continue
        tag = elem.tag
        concept = tag.split("}")[-1] if "}" in tag else tag
        if concept.lower() in ("context", "unit", "schemaref", "identifier", "period"):
            continue
        value = (elem.text or "").strip()
        if not value:
            continue
        # Format concept as readable predicate
        predicate = re.sub(r"([a-z])([A-Z])", r"\1_\2", concept).lower()
        unit_ref = elem.get("unitRef", "")
        unit_suffix = " (USD)" if unit_ref and "usd" in unit_ref.lower() else ""
        evidence = f"{concept}: {value}{unit_suffix}"
        candidates.append({
            "subject": entity_label,
            "predicate": predicate,
            "object": value + (unit_suffix.strip() or ""),
            "object_type": "literal:string",
            "tier": 0,
            "confidence": 1.0,
            "evidence": evidence,
            "locator": {"kind": "xml", "file_path": str(path)},
        })

    block_text = f"XBRL filing: {path.name}\n" + "\n".join(
        c["evidence"] for c in candidates[:50]
    )
    blocks = [DocumentBlock(
        text=block_text,
        locator={"kind": "xml", "file_path": str(path)},
    )]

    return ExtractedDocument(
        blocks=blocks,
        source_path=str(path),
        format="xbrl",
        tier0_candidates=candidates if candidates else None,
        metadata={"concept_count": len(candidates)},
    )


def extract_ical(path: Path) -> ExtractedDocument:
    """iCalendar — calendar events as tier-0 candidates."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    candidates = []
    current: dict = {}
    in_event = False

    for line in raw.splitlines():
        line = line.strip()
        if line == "BEGIN:VEVENT":
            in_event = True
            current = {}
        elif line == "END:VEVENT":
            summary = current.get("SUMMARY", "Untitled Event")
            dtstart = current.get("DTSTART", "")
            dtend = current.get("DTEND", "")
            location = current.get("LOCATION", "")
            evidence = f"{summary} from {dtstart} to {dtend}"
            if location:
                evidence += f" at {location}"
            candidates.append({
                "subject": summary,
                "predicate": "scheduled_at",
                "object": dtstart,
                "object_type": "literal:string",
                "tier": 0,
                "confidence": 1.0,
                "evidence": evidence,
                "locator": {"kind": "txt", "file_path": str(path)},
            })
            if location:
                candidates.append({
                    "subject": summary,
                    "predicate": "located_at",
                    "object": location,
                    "object_type": "literal:string",
                    "tier": 0,
                    "confidence": 1.0,
                    "evidence": evidence,
                    "locator": {"kind": "txt", "file_path": str(path)},
                })
            in_event = False
        elif in_event and ":" in line:
            key_part, _, val = line.partition(":")
            key = key_part.split(";")[0]
            current[key] = val

    block_text = "\n".join(c["evidence"] for c in candidates)
    blocks = [DocumentBlock(
        text=block_text or raw,
        locator={"kind": "txt", "file_path": str(path)},
    )]

    return ExtractedDocument(
        blocks=blocks,
        source_path=str(path),
        format="ical",
        tier0_candidates=candidates if candidates else None,
        metadata={"event_count": len(candidates)},
    )


def extract_rss(path: Path) -> ExtractedDocument:
    """RSS/Atom feeds — feed items as tier-0 candidates."""
    import xml.etree.ElementTree as ET

    raw = path.read_text(encoding="utf-8", errors="replace")
    candidates = []

    def _text(elem: Any, tag: str) -> str:
        child = elem.find(tag)
        return (child.text or "").strip() if child is not None else ""

    try:
        root = ET.fromstring(raw)
    except ET.ParseError:
        block = DocumentBlock(text=raw, locator={"kind": "xml", "file_path": str(path)})
        return ExtractedDocument(blocks=[block], source_path=str(path), format="rss")

    # RSS
    for item in root.iter("item"):
        title = _text(item, "title")
        if not title:
            continue
        link = _text(item, "link")
        pub_date = _text(item, "pubDate")
        desc = _text(item, "description")[:200]
        evidence = f"{title} ({pub_date}): {desc}" if desc else f"{title} ({pub_date})"
        candidates.append({
            "subject": title,
            "predicate": "published_at",
            "object": pub_date or link,
            "object_type": "literal:string",
            "tier": 0,
            "confidence": 1.0,
            "evidence": evidence,
            "locator": {"kind": "xml", "file_path": str(path)},
        })

    # Atom
    ATOM = "http://www.w3.org/2005/Atom"
    for entry in root.iter(f"{{{ATOM}}}entry"):
        title_elem = entry.find(f"{{{ATOM}}}title")
        title = (title_elem.text or "").strip() if title_elem is not None else ""
        if not title:
            continue
        updated_elem = entry.find(f"{{{ATOM}}}updated")
        updated = (updated_elem.text or "").strip() if updated_elem is not None else ""
        link_elem = entry.find(f"{{{ATOM}}}link")
        link = (link_elem.get("href") or "") if link_elem is not None else ""
        evidence = f"{title} ({updated})" if updated else title
        candidates.append({
            "subject": title,
            "predicate": "published_at",
            "object": updated or link,
            "object_type": "literal:string",
            "tier": 0,
            "confidence": 1.0,
            "evidence": evidence,
            "locator": {"kind": "xml", "file_path": str(path)},
        })

    block_text = "\n".join(c["evidence"] for c in candidates)
    blocks = [DocumentBlock(
        text=block_text or raw,
        locator={"kind": "xml", "file_path": str(path)},
    )]

    return ExtractedDocument(
        blocks=blocks,
        source_path=str(path),
        format="rss",
        tier0_candidates=candidates if candidates else None,
        metadata={"item_count": len(candidates)},
    )


# ---------------------------------------------------------------------------
# Router
# ---------------------------------------------------------------------------

_EXTRACTORS = {
    ".txt": extract_txt, ".text": extract_txt,
    ".md": extract_markdown, ".markdown": extract_markdown,
    ".html": extract_html, ".htm": extract_html,
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx, ".xls": extract_xlsx,
    ".pptx": extract_pptx,
    ".csv": extract_csv, ".tsv": extract_csv,
    ".json": extract_json, ".jsonl": extract_json,
    ".xml": extract_xml,
    ".xbrl": extract_xbrl,
    ".ics": extract_ical, ".ical": extract_ical,
    ".rss": extract_rss, ".atom": extract_rss,
}

SUPPORTED_EXTENSIONS = sorted(_EXTRACTORS.keys())


def extract(path: Path) -> ExtractedDocument:
    """Extract content from any supported format.

    Returns ExtractedDocument with:
    - blocks: List[DocumentBlock] for unstructured content (LLM path)
    - tier0_candidates: List[dict] for structured data (bypass path)
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")
    ext = path.suffix.lower()
    if ext not in _EXTRACTORS:
        raise ValueError(f"Unsupported format: {ext}\nSupported: {', '.join(SUPPORTED_EXTENSIONS)}")
    return _EXTRACTORS[ext](path)
